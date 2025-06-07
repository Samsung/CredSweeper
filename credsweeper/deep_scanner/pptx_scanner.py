import io
import logging
from abc import ABC
from typing import List, Optional

from pptx import Presentation

from credsweeper.credentials.candidate import Candidate
from credsweeper.deep_scanner.abstract_scanner import AbstractScanner
from credsweeper.file_handler.data_content_provider import DataContentProvider
from credsweeper.file_handler.string_content_provider import StringContentProvider

logger = logging.getLogger(__name__)


class PptxScanner(AbstractScanner, ABC):
    """Implements pptx scanning"""

    def data_scan(
            self,  #
            data_provider: DataContentProvider,  #
            depth: int,  #
            recursive_limit_size: int) -> Optional[List[Candidate]]:
        """Tries to scan pptx text elements for all slides"""
        try:
            candidates = []
            pptx_lines = []
            presentation = Presentation(io.BytesIO(data_provider.data))
            for n, slide in enumerate(presentation.slides):
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            pptx_lines.append(paragraph.text)
                string_data_provider = StringContentProvider(lines=pptx_lines,
                                                             file_path=data_provider.file_path,
                                                             file_type=data_provider.file_type,
                                                             info=f"{data_provider.info}|PPTX:{n+1}")
                pptx_candidates = self.scanner.scan(string_data_provider)
                candidates.extend(pptx_candidates)
            return candidates
        except Exception as pptx_exc:
            logger.error(f"{data_provider.file_path}:{pptx_exc}")
        return None
